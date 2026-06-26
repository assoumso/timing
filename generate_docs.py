import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

from reportlab.lib.pagesizes import letter
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib import colors

# Colors matching the application
COLOR_NAVY = RGBColor(11, 73, 152)      # #0b4998
COLOR_ORANGE = RGBColor(238, 123, 17)   # #ee7b11
COLOR_DARK = RGBColor(30, 41, 59)       # #1e293b
COLOR_WHITE = RGBColor(255, 255, 255)
COLOR_LIGHT_BG = RGBColor(248, 250, 252)

def set_slide_bg(slide, rgb_color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = rgb_color

def add_header(slide, title_text):
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.33), Inches(0.8))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = COLOR_NAVY
    p.font.name = 'Arial'

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Slide 1: Title
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, COLOR_NAVY)
    txBox = slide.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11.33), Inches(3.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "ERP BARAKAT"
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = COLOR_ORANGE
    p.font.name = 'Arial'
    p2 = tf.add_paragraph()
    p2.text = "Guide Procédural & Tutoriels Étape par Étape"
    p2.font.size = Pt(28)
    p2.font.bold = True
    p2.font.color.rgb = COLOR_WHITE
    p2.font.name = 'Arial'
    
    def add_bullet(tf, bold_prefix, text_content, level=0):
        p = tf.add_paragraph()
        p.level = level
        run1 = p.add_run()
        run1.text = bold_prefix
        run1.font.bold = True
        run1.font.size = Pt(15)
        run1.font.color.rgb = COLOR_DARK
        run2 = p.add_run()
        run2.text = text_content
        run2.font.size = Pt(14)
        run2.font.color.rgb = COLOR_DARK

    # Slide 2: Création Utilisateur
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, COLOR_LIGHT_BG)
    add_header(slide, "Tutoriel 1 : Créer un utilisateur & attribuer des droits")
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(12.33), Inches(5.8))
    tf = txBox.text_frame
    tf.word_wrap = True
    add_bullet(tf, "1. Ouvrir l'onglet : ", "Sélectionnez l'onglet 'Administration' dans le menu principal.")
    add_bullet(tf, "2. Remplir le formulaire : ", "Allez à 'Créer un nouveau compte utilisateur', entrez le Nom complet, l'Email/Identifiant, et le Mot de passe.")
    add_bullet(tf, "3. Sélectionner le rôle : ", "Choisissez le rôle approprié dans la liste déroulante (Super Admin, Directeur, Informaticien, Professeur).")
    add_bullet(tf, "4. Enregistrer : ", "Cliquez sur le bouton orange 'Créer le compte utilisateur'.")
    add_bullet(tf, "5. Assigner la Saisie des Moyennes : ", "Dans la liste des comptes ci-dessous, repérez l'utilisateur et cochez la case 'Saisie des Moyennes'.")

    # Slide 3: Planifier Cours
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, COLOR_LIGHT_BG)
    add_header(slide, "Tutoriel 2 : Planifier un cours sur l'Emploi du Temps")
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(12.33), Inches(5.8))
    tf = txBox.text_frame
    tf.word_wrap = True
    add_bullet(tf, "1. Accéder au module : ", "Allez sur l'onglet 'EDT - Planification'.")
    add_bullet(tf, "2. Remplir les détails : ", "Sélectionnez la Classe, le Professeur, la Matière et la Salle d'étude.")
    add_bullet(tf, "3. Choisir le créneau : ", "Choisissez le Jour de la semaine (ex: Lundi) et le Créneau horaire (ex: 08:00 - 10:00).")
    add_bullet(tf, "4. Détection automatique : ", "Le système bloque instantanément la planification si le prof ou la salle est déjà occupé(e).")
    add_bullet(tf, "5. Validation : ", "Cliquez sur 'Planifier ce cours' pour l'ajouter à la grille hebdomadaire.")

    # Slide 4: Saisie des Notes & Figeage
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, COLOR_LIGHT_BG)
    add_header(slide, "Tutoriel 3 : Saisie des notes (Enseignants) & Verrouillage")
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(12.33), Inches(5.8))
    tf = txBox.text_frame
    tf.word_wrap = True
    add_bullet(tf, "1. Connexion Professeur : ", "Connectez-vous au portail avec vos identifiants d'enseignant.")
    add_bullet(tf, "2. Sélectionner la classe : ", "Allez dans 'Espace Enseignant' > 'Saisie des Moyennes', puis choisissez la Classe et la Matière.")
    add_bullet(tf, "3. Saisir les notes : ", "Entrez la note de chaque élève de la liste (valeur sur 20) et cliquez sur 'Enregistrer'.")
    add_bullet(tf, "4. Figer définitivement : ", "Cliquez sur le bouton vert '🔒 Valider & Figer' pour verrouiller la saisie.")
    add_bullet(tf, "5. Sécurité : ", "Une fois figée, la note affiche le badge '🔒 Figé' et le professeur ne peut plus la modifier.")

    # Slide 5: Déverrouillage Direction
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, COLOR_LIGHT_BG)
    add_header(slide, "Tutoriel 4 : Déverrouiller une note (Direction)")
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(12.33), Inches(5.8))
    tf = txBox.text_frame
    tf.word_wrap = True
    add_bullet(tf, "1. Accéder au portail prof : ", "Connectez-vous en tant que Directeur ou Super Admin et allez dans le portail enseignant.")
    add_bullet(tf, "2. Trouver la note figée : ", "Sélectionnez la classe et la matière concernée dans le filtre du portail.")
    add_bullet(tf, "3. Déverrouiller : ", "Cliquez sur le lien orange '🔓 Déverr.' à côté de la note verrouillée.")
    add_bullet(tf, "4. Confirmation : ", "Confirmez l'alerte du navigateur. Le badge '🔒 Figé' disparaît et l'enseignant peut à nouveau corriger.")

    # Slide 6: Validation PV
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, COLOR_LIGHT_BG)
    add_header(slide, "Tutoriel 5 : Valider le Procès-Verbal (PV)")
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(12.33), Inches(5.8))
    tf = txBox.text_frame
    tf.word_wrap = True
    add_bullet(tf, "1. Accéder à l'onglet : ", "Allez sur 'Évaluations & Examens' > sous-onglet 'Procès-Verbaux & Validation'.")
    add_bullet(tf, "2. Vérifier la classe : ", "Sélectionnez la classe et examinez les moyennes affichées sur l'écran (Natte de vérification).")
    add_bullet(tf, "3. Lancer la validation : ", "Cliquez sur 'Valider le PV pour cette Classe' pour verrouiller le trimestre.")
    add_bullet(tf, "4. Effet immédiat : ", "Les bulletins scolaires de cette classe deviennent instantanément disponibles pour l'impression.")

    # Slide 7: Scolarité & Facturation
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, COLOR_LIGHT_BG)
    add_header(slide, "Tutoriel 6 : Enregistrer un paiement & générer le reçu")
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(12.33), Inches(5.8))
    tf = txBox.text_frame
    tf.word_wrap = True
    add_bullet(tf, "1. Ouvrir l'onglet financier : ", "Sélectionnez l'onglet 'Finances & Écolages'.")
    add_bullet(tf, "2. Sélectionner l'élève : ", "Recherchez et sélectionnez l'élève dans la liste déroulante.")
    add_bullet(tf, "3. Remplir le versement : ", "Entrez le Montant payé, le Motif du versement, et le Mode de règlement (Espèces, Chèque, etc.).")
    add_bullet(tf, "4. Enregistrer : ", "Cliquez sur 'Enregistrer la transaction'. Le solde de l'élève se met à jour immédiatement.")
    add_bullet(tf, "5. Télécharger/Envoyer : ", "Cliquez sur 'Générer Reçu PDF' pour l'imprimer, ou utilisez le bouton WhatsApp pour envoyer la notification.")

    prs.save('guide_utilisation.pptx')
    print("PowerPoint generated successfully.")

def create_markdown_guide():
    content = """# 📘 GUIDE D'UTILISATION PROCÉDURAL COMPLET (PAS-À-PAS)
## ERP Scolaire & Emploi du Temps (« BARAKAT »)

Ce guide détaille la marche à suivre pas-à-pas pour exécuter chaque tâche sur l'application.

---

## 🗺️ Index des Procédures
1. [Comment créer un utilisateur et configurer ses droits (RBAC)](#1-comment-créer-un-utilisateur-et-configurer-ses-droits-rbac)
2. [Comment configurer les matières, coefficients et langues LV2](#2-comment-configurer-les-matières-coefficients-et-langues-lv2)
3. [Comment inscrire un nouvel élève](#3-comment-inscrire-un-nouvel-élève)
4. [Comment affecter un professeur principal à une classe](#4-comment-affecter-un-professeur-principal-à-une-classe)
5. [Comment planifier un cours sur l'Emploi du Temps (EDT)](#5-comment-planifier-un-cours-sur-lemploi-du-temps-edt)
6. [Comment déclarer les indisponibilités d'un enseignant](#6-comment-déclarer-les-indisponibilités-dun-enseignant)
7. [Comment saisir les notes (côté Enseignant) et figer la saisie](#7-comment-saisir-les-notes-côté-enseignant-et-figer-la-saisie)
8. [Comment déverrouiller une note validée par un enseignant](#8-comment-déverrouiller-une-note-validée-par-un-enseignant)
9. [Comment saisir ou supprimer des notes (côté Administration)](#9-comment-saisir-ou-supprimer-des-notes-côté-administration)
10. [Comment valider le procès-verbal (PV) d'une classe](#10-comment-valider-le-procès-verbal-pv-dune-classe)
11. [Comment imprimer ou exporter les bulletins de notes](#11-comment-imprimer-ou-exporter-les-bulletins-de-notes)
12. [Comment enregistrer un paiement d'écolage et envoyer le reçu PDF](#12-comment-enregistrer-un-paiement-décolage-et-envoyer-le-reçu-pdf)

---

### 1. Comment créer un utilisateur et configurer ses droits (RBAC)
*   **Étape 1 :** Cliquez sur l'onglet **Administration** dans la barre de navigation supérieure.
*   **Étape 2 :** Dans la section **"Créer un nouveau compte utilisateur"**, remplissez les champs :
    *   **Nom complet** (ex: *Koffi Yao*)
    *   **Identifiant / Email** (ex: *koffi.yao@ecole.ci*)
    *   **Mot de passe**
*   **Étape 3 :** Choisissez le **Rôle** dans le menu déroulant :
    *   `Super Admin` ou `Directeur` (accès complet)
    *   `Informaticien / Correspondant Fichier` (accès technique)
    *   `Professeur` (accès enseignant restreint)
*   **Étape 4 :** Cliquez sur le bouton orange **"Créer le compte utilisateur"**.
*   **Étape 5 (Permissions) :** Dans la table **"Comptes utilisateurs enregistrés"** située juste en dessous, recherchez l'utilisateur créé.
*   **Étape 6 :** Cochez la case **"Saisie des Moyennes"** (`saisie_moyennes`) pour l'autoriser à saisir les notes de ses classes. La permission prend effet instantanément.

---

### 2. Comment configurer les matières, coefficients et langues LV2
*   **Étape 1 :** Accédez à l'onglet **Matières** du menu.
*   **Étape 2 :** Pour ajouter une matière, allez au formulaire **"Ajouter une Matière"** :
    *   Saisissez le **Nom de la matière** (ex: *Allemand*).
    *   Sélectionnez le **Coefficient** (poids pour le calcul des moyennes).
    *   Cochez **"Langue Vivante 2 (LV2)"** s'il s'agit d'une matière facultative ou de langue secondaire.
    *   Sélectionnez la couleur thématique pour l'affichage dans l'Emploi du Temps.
*   **Étape 3 :** Cliquez sur le bouton orange **"Ajouter la Matière"**.

---

### 3. Comment inscrire un nouvel élève
*   **Étape 1 :** Cliquez sur l'onglet **Élèves**.
*   **Étape 2 :** Dans le panneau **"Inscrire un nouvel élève"**, complétez les informations requises :
    *   **Nom** et **Prénom** de l'élève.
    *   **Genre** (Masculin / Féminin).
    *   **Classe d'affectation** (sélectionnez dans la liste déroulante des classes existantes).
*   **Étape 3 :** Cliquez sur **"Enregistrer l'inscription"**. L'élève apparaît immédiatement dans le registre général de sa classe.

---

### 4. Comment affecter un professeur principal à une classe
*   **Étape 1 :** Rendez-vous dans l'onglet **EDT - Salles & Classes**.
*   **Étape 2 :** Dans la table des classes enregistrées, localisez la classe concernée et cliquez sur son bouton de modification (ou modifiez directement les paramètres de création de la classe).
*   **Étape 3 :** Dans le menu déroulant **"Professeur Principal"**, sélectionnez l'enseignant référent de cette classe.
*   **Étape 4 :** Cliquez sur **"Enregistrer les modifications"**. Le nom du professeur principal sera automatiquement imprimé sur tous les bulletins scolaires de cette classe.

---

### 5. Comment planifier un cours sur l'Emploi du Temps (EDT)
*   **Étape 1 :** Allez sur l'onglet **EDT - Planification**.
*   **Étape 2 :** Dans le formulaire de gauche, sélectionnez :
    *   La **Classe** cible.
    *   Le **Professeur** en charge.
    *   La **Matière** enseignée.
    *   La **Salle** réservée.
*   **Étape 3 :** Indiquez le **Jour** de la semaine et le **Créneau Horaire** (ex: *Mardi de 10h à 12h*).
*   **Étape 4 :** Cliquez sur **"Planifier ce cours"**.
    *   *Si un conflit est détecté (enseignant occupé, salle déjà prise ou professeur indisponible), une boîte d'alerte rouge bloque la saisie.*
    *   *Si tout est correct, le cours s'inscrit automatiquement dans la grille hebdomadaire visible à l'écran.*

---

### 6. Comment déclarer les indisponibilités d'un enseignant
*   **Étape 1 :** Rendez-vous dans l'**Espace Enseignant** (Portail Professeur).
*   **Étape 2 :** Sélectionnez le professeur concerné dans la liste déroulante supérieure.
*   **Étape 3 :** Repérez la grille horaire de disponibilité.
*   **Étape 4 :** Cliquez directement sur les cases de la grille pour basculer leur état :
    *   Une case grisée avec le symbole **❌ Indisponible** signifie que le professeur ne peut pas enseigner sur ce créneau.
    *   Les modifications se sauvegardent automatiquement dans la base locale.

---

### 7. Comment saisir les notes (côté Enseignant) et figer la saisie
*   **Étape 1 :** Ouvrez le **Portail Enseignant** > section **Saisie des Moyennes**.
*   **Étape 2 :** Sélectionnez la **Classe** et la **Discipline (Matière)** dans les filtres en haut.
*   **Étape 3 :** Dans la liste des élèves qui s'affiche, entrez la note sur 20 dans le champ numérique de chaque élève.
*   **Étape 4 :** Saisissez l'intitulé de l'évaluation (ex: *Interrogation n°1*) et son coefficient.
*   **Étape 5 :** Cliquez sur le bouton orange **"Enregistrer toutes les notes saisies"**.
*   **Étape 6 (Verrouillage) :** Pour figer définitivement ces moyennes, repérez la section **"Notes déjà saisies"** à droite et cliquez sur le bouton vert **"🔒 Valider & Figer"** :
    *   Les notes affichent désormais le badge **🔒 Figé**.
    *   Le professeur ne peut plus modifier ni supprimer ces notes.

---

### 8. Comment déverrouiller une note validée par un enseignant
*   **Étape 1 :** Connectez-vous avec un compte **Directeur** ou **Super Admin**.
*   **Étape 2 :** Accédez au **Portail Enseignant** (Saisie des Moyennes) et sélectionnez la classe et la matière de la note à corriger.
*   **Étape 3 :** Dans la liste des notes saisies à droite, localisez la note portant le badge **🔒 Figé**.
*   **Étape 4 :** Cliquez sur le bouton bleu **"🔓 Déverr."** situé à côté de la note.
*   **Étape 5 :** Confirmez la demande dans la boîte de dialogue. Le badge disparaît et l'enseignant peut de nouveau modifier la note depuis sa session.

---

### 9. Comment saisir ou supprimer des notes (côté Administration)
*   **Étape 1 :** Allez sur l'onglet **Évaluations & Examens**.
*   **Étape 2 :** *Si vous êtes autorisé*, le formulaire de saisie de gauche s'affiche.
    *   Sélectionnez l'Élève, la Matière, le Coefficient, le type de devoir et la note.
    *   Cliquez sur **"Soumettre au Registre Électronique"**.
*   **Étape 3 (Suppression) :** Pour supprimer une note erronée, repérez-la dans la liste de droite **"Derniers coefficients & notes enregistrés"**.
*   **Étape 4 :** Cliquez sur le bouton rouge **Corbeille (Supprimer)** et validez la confirmation.

---

### 10. Comment valider le procès-verbal (PV) d'une classe
*   **Étape 1 :** Rendez-vous dans l'onglet **Évaluations & Examens**, sous-onglet **Procès-Verbaux & Validation**.
*   **Étape 2 :** Sélectionnez la **Classe** souhaitée.
*   **Étape 3 :** Examinez le tableau de la "Natte de vérification". Assurez-vous que toutes les moyennes des matières sont présentes et valides.
*   **Étape 4 :** Cliquez sur le bouton bleu **"Valider le PV pour cette Classe"**. Le PV passe au statut **"Validé"** et fige les moyennes trimestrielles de la classe.

---

### 11. Comment imprimer ou exporter les bulletins de notes
*   **Étape 1 :** Assurez-vous d'abord que le PV de la classe a été validé (voir rubrique précédente).
*   **Étape 2 :** Allez sur l'onglet **Évaluations & Examens** > sous-onglet **Éditeur de Bulletins**.
*   **Étape 3 :** Sélectionnez la **Classe** et l'**Élève**.
    *   *Si le PV est validé, le bulletin s'affiche avec la moyenne générale, le rang de l'élève, et l'appréciation du professeur principal.*
    *   *Si le PV n'est pas validé, un écran d'avertissement jaune bloque la consultation.*
*   **Étape 4 :** Cliquez sur le bouton bleu **"Imprimer le Bulletin"** pour ouvrir l'aperçu avant impression système ou le sauvegarder en PDF.

---

### 12. Comment enregistrer un paiement d'écolage et envoyer le reçu PDF
*   **Étape 1 :** Cliquez sur l'onglet **Finances & Écolages** dans le menu principal.
*   **Étape 2 :** Recherchez et sélectionnez l'élève dans la liste déroulante.
*   **Étape 3 :** Remplissez les champs de paiement de la section **"Enregistrer une nouvelle transaction"** :
    *   **Montant versé** (ex: *50 000 FCFA*)
    *   **Motif du paiement** (ex: *Scolarité 1er Trimestre*)
    *   **Mode de règlement** (Espèces, Orange Money, Wave, etc.)
*   **Étape 4 :** Cliquez sur **"Enregistrer la transaction"**. Le solde restant de l'élève est recalculé.
*   **Étape 5 :** Dans la table des reçus en bas de page, cliquez sur **"Reçu PDF"** pour télécharger le justificatif.
*   **Étape 6 :** Cliquez sur le bouton vert **"WhatsApp / SMS"** pour générer le lien de partage rapide et l'envoyer au tuteur de l'élève.
"""
    with open("GUIDE_UTILISATION.md", "w", encoding="utf-8") as f:
        f.write(content)
    print("Markdown generated successfully.")

def create_pdf_guide():
    doc = SimpleDocTemplate("guide_utilisation.pdf", pagesize=letter,
                            rightMargin=30, leftMargin=30, topMargin=30, bottomMargin=30)
    styles = getSampleStyleSheet()
    
    style_title = ParagraphStyle(
        name='TitleStyle',
        parent=styles['Heading1'],
        fontName='Helvetica-Bold',
        fontSize=20,
        textColor=colors.HexColor('#0b4998'),
        spaceAfter=10
    )
    style_subtitle = ParagraphStyle(
        name='SubTitleStyle',
        parent=styles['Normal'],
        fontName='Helvetica-Bold',
        fontSize=12,
        textColor=colors.HexColor('#ee7b11'),
        spaceAfter=20
    )
    style_h2 = ParagraphStyle(
        name='H2Style',
        parent=styles['Heading2'],
        fontName='Helvetica-Bold',
        fontSize=12,
        textColor=colors.HexColor('#0b4998'),
        spaceBefore=12,
        spaceAfter=6
    )
    style_body = ParagraphStyle(
        name='BodyStyle',
        parent=styles['Normal'],
        fontName='Helvetica',
        fontSize=9,
        textColor=colors.HexColor('#1e293b'),
        spaceBefore=4,
        spaceAfter=4,
        leading=12
    )

    story = []
    story.append(Paragraph("GUIDE D'UTILISATION PAS-À-PAS - ERP BARAKAT", style_title))
    story.append(Paragraph("Manuel détaillé des procédures d'utilisation de l'établissement", style_subtitle))
    
    procedures = [
        ("1. Créer un utilisateur et configurer ses droits (RBAC)", 
         "• Allez sur l'onglet Administration.\n"
         "• Dans le formulaire 'Créer un nouveau compte', saisissez Nom complet, Identifiant/Email et Mot de passe.\n"
         "• Choisissez le Rôle (Super Admin, Directeur, Informaticien, Professeur).\n"
         "• Cliquez sur 'Créer le compte utilisateur'.\n"
         "• Cochez la case 'Saisie des Moyennes' sous le profil créé pour activer son droit de saisie."),
         
        ("2. Configurer les matières, coefficients et langues LV2", 
         "• Allez sur l'onglet Matières.\n"
         "• Entrez le nom de la matière, sélectionnez son Coefficient, et cochez 'LV2' si c'est une langue facultative.\n"
         "• Choisissez une couleur d'EDT et cliquez sur 'Ajouter la Matière'."),
         
        ("3. Inscrire un nouvel élève", 
         "• Allez sur l'onglet Élèves.\n"
         "• Saisissez Nom, Prénom, Genre et sélectionnez la classe d'affectation.\n"
         "• Cliquez sur 'Enregistrer l'inscription'."),
         
        ("4. Planifier un cours sur l'Emploi du Temps (EDT)", 
         "• Allez sur l'onglet EDT - Planification.\n"
         "• Sélectionnez la Classe, le Professeur, la Matière, la Salle, le Jour et le Créneau horaire.\n"
         "• Cliquez sur 'Planifier ce cours'. Le système rejette le cours si un conflit de salle, de professeur ou d'indisponibilité est détecté."),
         
        ("5. Saisir les notes (Enseignants) et figer la saisie", 
         "• Ouvrez l'Espace Enseignant > Saisie des Moyennes.\n"
         "• Sélectionnez la classe et la matière.\n"
         "• Entrez les notes sur 20, le motif, le coefficient et cliquez sur 'Enregistrer'.\n"
         "• Pour figer la saisie et bloquer les modifications, cliquez sur 'Valider & Figer' dans la liste des notes saisies à droite."),
         
        ("6. Déverrouiller une note validée", 
         "• Connectez-vous avec un compte Directeur ou Super Admin.\n"
         "• Allez sur le Portail Enseignant > Saisie des Moyennes et filtrez par classe/matière.\n"
         "• Cliquez sur le lien 'Déverr.' à côté de la note verrouillée pour autoriser à nouveau les modifications."),
         
        ("7. Valider le Procès-Verbal (PV) de classe", 
         "• Allez sur l'onglet Évaluations > Procès-Verbaux & Validation.\n"
         "• Sélectionnez la classe, vérifiez la natte de notes.\n"
         "• Cliquez sur 'Valider le PV pour cette Classe' pour autoriser l'édition des bulletins."),
         
        ("8. Imprimer les bulletins de notes", 
         "• Allez sur l'onglet Évaluations > Éditeur de Bulletins (disponible uniquement après validation du PV de la classe).\n"
         "• Sélectionnez la classe, l'élève, et cliquez sur 'Imprimer le Bulletin'."),
         
        ("9. Enregistrer un paiement et éditer le reçu", 
         "• Allez sur l'onglet Finances & Ecolages.\n"
         "• Sélectionnez l'élève, entrez le montant, le motif, le mode de règlement, et cliquez sur 'Enregistrer'.\n"
         "• Cliquez sur 'Reçu PDF' pour le télécharger ou sur le bouton WhatsApp pour envoyer la notification.")
    ]

    for title, steps in procedures:
        story.append(Paragraph(title, style_h2))
        for step in steps.split('\n'):
            story.append(Paragraph(step, style_body))
        story.append(Spacer(1, 6))

    doc.build(story)
    print("PDF generated successfully.")

if __name__ == "__main__":
    create_presentation()
    create_markdown_guide()
    create_pdf_guide()
